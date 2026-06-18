import os
import json
import pandas as pd
 
 
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read().strip()
 
 
def read_pdf(file_path):
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(file_path)
        if not text or not text.strip():
            return "No readable text found in PDF."
        return text.strip()
    except Exception as e:
        return f"Failed to read PDF: {str(e)}"
 
 
def read_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        return (
            f"CSV File Analysis\n\n"
            f"Columns: {', '.join(df.columns)}\n\n"
            f"Rows: {len(df)}\n\n"
            f"Sample Data:\n"
            f"{df.head(10).to_string()}"
        )
    except Exception as e:
        return f"Failed to read CSV: {str(e)}"
 
 
def read_excel(file_path):
    try:
        excel_file = pd.ExcelFile(file_path)
        result = []
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet)
            result.append(f"\nSheet Name: {sheet}\n")
            result.append(f"Columns: {', '.join(str(c) for c in df.columns)}\n")
            result.append(f"Rows: {len(df)}\n")
            result.append("Sample Data:\n")
            result.append(df.head(10).to_string())
            result.append("\n")
        return "\n".join(result)
    except Exception as e:
        return f"Failed to read Excel file: {str(e)}"
 
 
def read_docx(file_path):
    try:
        import docx
        doc     = docx.Document(file_path)
        content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
 
        # Tables in .docx are common for budgets/specs/comparisons and are
        # otherwise silently dropped — pull their text too.
        for ti, table in enumerate(doc.tables, start=1):
            rows_text = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    rows_text.append(" | ".join(cells))
            if rows_text:
                content += f"\n\n[Table {ti}]\n" + "\n".join(rows_text)
 
        return content or "No readable text found in Word document."
    except Exception as e:
        return f"Failed to read Word document: {str(e)}"
 
 
def read_pptx(file_path):
    """Extract text from all slides (titles, body text, tables, notes)."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
 
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_text.append(" | ".join(cells))
 
            # Speaker notes often carry context/data not shown on the slide
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[Notes] {notes}")
 
            if slide_text:
                parts.append(f"\nSlide {i}:\n" + "\n".join(slide_text))
 
        content = "\n".join(parts)
        return content or "No readable text found in PowerPoint file."
    except Exception as e:
        return f"Failed to read PowerPoint file: {str(e)}"
 
 
def read_json(file_path):
    """
    If the JSON is an array of records (or {key: [records]}), summarize it
    like a CSV (columns, row count, sample rows) so downstream analysis can
    treat it as tabular data. Otherwise, return the pretty-printed structure
    as text.
    """
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            raw = json.load(f)
    except Exception as e:
        return f"Failed to read JSON file: {str(e)}"
 
    # Try to find a list-of-dicts to treat as tabular, either at the top
    # level or as the first list-valued key in a top-level dict.
    records = None
    if isinstance(raw, list) and raw and isinstance(raw[0], dict):
        records = raw
    elif isinstance(raw, dict):
        for v in raw.values():
            if isinstance(v, list) and v and isinstance(v[0], dict):
                records = v
                break
 
    if records is not None:
        try:
            df = pd.json_normalize(records)
            return (
                f"JSON File Analysis (tabular)\n\n"
                f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
                f"Rows: {len(df)}\n\n"
                f"Sample Data:\n"
                f"{df.head(10).to_string()}"
            )
        except Exception:
            pass  # fall through to raw dump
 
    pretty = json.dumps(raw, indent=2, default=str)
    return f"JSON File Content\n\n{pretty[:8000]}"
 
 
def parse_file_only(file_path, file_ext):
    ext = file_ext.lower()
    if ext == ".txt" or ext == ".md":
        return read_txt(file_path)
    elif ext == ".pdf":
        return read_pdf(file_path)
    elif ext == ".csv":
        return read_csv(file_path)
    elif ext in (".xlsx", ".xls"):
        return read_excel(file_path)
    elif ext in (".docx", ".doc"):
        return read_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx(file_path)
    elif ext == ".json":
        return read_json(file_path)
    else:
        # Last resort: try reading as text
        try:
            return read_txt(file_path)
        except Exception:
            raise ValueError(f"Unsupported file extension: {ext}")
 
 
def parse_multiple_files(file_paths: list) -> str:
    """Reads and combines content from multiple files into one string."""
    if not file_paths:
        return ""
    combined_parts = []
    for idx, file_path in enumerate(file_paths, start=1):
        filename = os.path.basename(file_path)
        ext      = os.path.splitext(file_path)[1].lower()
        try:
            content = parse_file_only(file_path, ext)
        except Exception as e:
            content = f"Could not read file: {str(e)}"
        combined_parts.append(
            f"--- File {idx}: {filename} ---\n{content}\n--- End of File {idx} ---"
        )
    return "\n\n".join(combined_parts)
 
 
def get_user_input():
    print("Choose input type:")
    print("1. Type prompt manually")
    print("2. Load from file")
    choice = input("Enter choice: ").strip()
 
    if choice == "1":
        return input("Enter your prompt:\n>>> ").strip()
 
    elif choice == "2":
        instruction = input("Enter instruction:\n>>> ").strip()
        file_path   = input("Enter file path:\n>>> ").strip().strip('"')
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        ext     = os.path.splitext(file_path)[1].lower()
        content = parse_file_only(file_path, ext)
        return {"prompt": instruction, "content": content}
    else:
        raise ValueError("Invalid choice.")
 