import os
import re
import textwrap
import logging
from datetime import datetime
 
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
 
logger = logging.getLogger(__name__)
 
# ── Dimensions (16:9) ─────────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
 
TITLE_H  = Inches(1.1)
BODY_L   = Inches(0.45)
BODY_T   = Inches(1.35)
BODY_W   = SW - Inches(0.9)
BODY_H   = SH - Inches(1.35) - Inches(0.55)
FOOTER_T = SH - Inches(0.42)
 
# Colors
HP_BLUE  = RGBColor(0,   32,  91)
HP_RED   = RGBColor(228, 0,   43)
WHITE    = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(244, 246, 252)
DARK_TXT = RGBColor(20,  20,  20)
GREY_TXT = RGBColor(100, 100, 120)
 
# Height estimation constants
_BULLET_PAD = 10
_TOP_PAD    = 14
MIN_BG_H    = Inches(1.6)
 
MAX_BULLETS      = 8
MAX_SLIDES       = 14   # raised from 12 to match new 10–14 target
MAX_CHART_SLIDES = 4
 
FOOTER_MAX_CHARS = 42
 
ABBREVS = {'HPCL','BPCL','IOCL','LPG','ATF','KL','MT','OMC','HP','INR'}
 
ENTITY_KEYWORDS = {
    'hpcl', 'bpcl', 'iocl', 'reliance', 'nayara',
    'petrol', 'diesel', 'lpg', 'atf', 'lubricants', 'omc',
}
 
_FILTERED_HEADINGS = re.compile(
    r'^(title\s*page|acknowledgement|acknowledgment|thank\s*you|closing|'
    r'conclusion|introduction\s+to\s+.+)$',
    re.IGNORECASE
)
 
_seen_headings: set = set()
 
 
# ── Public entry point ────────────────────────────────────────────────────────
 
def generate_ppt(content, output_path="outputs/output.pptx",
                 references=None, filename_title="Untitled",
                 chart_images=None, add_acknowledgement=False):
    global _seen_headings
    _seen_headings = set()
 
    chart_images = list(chart_images or [])
    os.makedirs(
        os.path.dirname(output_path) if os.path.dirname(output_path) else '.',
        exist_ok=True
    )
 
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
 
    # 1. Cover
    _cover(prs, filename_title)
 
    # 2. Parse & cap content slides
    sections = _parse(content)[:MAX_SLIDES]
 
    placed = set()
    for sec_title, bullets, slide_type in sections:
        if not bullets:
            continue
 
        sec_text = sec_title + ' ' + ' '.join(_flatten(b) for b in bullets)
        sec_kw   = _kw(sec_text)
 
        best_ci, best_score = None, 0
        for ci, (ct, img_path) in enumerate(chart_images):
            if ci in placed or not os.path.exists(img_path):
                continue
            ckw = _kw(_clean(ct))
            score = len(sec_kw & ckw)
            score += 2 * len(sec_kw & ckw & ENTITY_KEYWORDS)
            if score > best_score:
                best_score, best_ci = score, ci
 
        sl = _blank(prs)
        _title_bar(sl, sec_title)
        _footer(sl, filename_title)
 
        # ── Slide-type-aware layout ───────────────────────────────────────────
        # The LLM marks special slides in the parsed output; we use those
        # signals to choose the right visual layout.
        if slide_type == "executive_summary":
            _add_executive_summary_slide(sl, bullets)
        elif slide_type == "kpi_dashboard":
            _add_kpi_dashboard_slide(sl, bullets)
        elif best_ci is not None and best_score > 0:
            _add_bullets_left(sl, bullets)
            _add_chart_right(sl, chart_images[best_ci][1])
            placed.add(best_ci)
        else:
            _add_bullets_full(sl, bullets)
 
    # 3. Remaining chart slides
    unplaced = [(ct, p) for ci, (ct, p) in enumerate(chart_images)
                if ci not in placed and os.path.exists(p)]
 
    pairs = _pair_charts(unplaced)[:MAX_CHART_SLIDES]
    for pair in pairs:
        sl = _blank(prs)
        _footer(sl, filename_title)
 
        if len(pair) == 2:
            ct, img_path = pair
            _title_bar(sl, _clean(ct))
            _add_chart_full(sl, img_path)
        else:
            (ct1, img1), (ct2, img2), score = pair
            if score > 0:
                _title_bar(sl, _chart_pair_title(ct1, ct2))
            else:
                _title_bar(sl, "Additional Data Visualisations")
            _add_two_charts(sl, img1, img2)
 
    # 4. Acknowledgement
    if add_acknowledgement:
        _acknowledgement(prs)
 
    # 5. Thank you — always last
    _thank_you(prs, filename_title)
 
    prs.save(output_path)
    total = len(prs.slides)
    logger.info(f"PPT saved: {output_path} ({total} slides)")
    print(f"PPT saved to: {output_path} ({total} slides)")
    return total
 
 
# ── Slide-type-aware layouts ──────────────────────────────────────────────────
 
def _add_executive_summary_slide(sl, bullets):
    """
    Executive Summary gets a special two-column layout:
    left column = 3 answer bullets (What / Why / Action),
    right column = coloured callout box for the headline insight.
    Falls back to full-width if fewer than 2 bullets.
    """
    clean = [b for b in (_clean_bullet(x) for x in bullets[:MAX_BULLETS]) if b]
    if not clean or len(clean) < 2:
        _add_bullets_full(sl, bullets)
        return
 
    # Headline callout box (first bullet — the most important insight)
    headline = clean[0]
    box_w = SW * 0.38
    box_h = Inches(2.2)
    box_t = BODY_T + Inches(0.3)
    box_l = SW - box_w - Inches(0.3)
 
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               box_l, box_t, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = HP_BLUE
    box.line.fill.background()
    _no_shadow(box)
 
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.2)
    tf.margin_top   = Inches(0.15)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
 
    # Label above the box
    tb_lbl = sl.shapes.add_textbox(box_l, BODY_T, box_w, Inches(0.35))
    p_lbl  = tb_lbl.text_frame.paragraphs[0]
    p_lbl.text = "HEADLINE INSIGHT"
    p_lbl.font.size = Pt(9)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = HP_RED
 
    # Remaining bullets on the left
    left_w = SW * 0.56 - Inches(0.1)
    remaining = clean[1:]
    if remaining:
        n      = len(remaining)
        avg_ch = max((len(b) for b in remaining), default=40)
        fsize  = _pick_fsize(n, avg_ch, wide=False)
        w_in   = float((left_w - Inches(0.3)) / 914400)
        bg_h   = _estimate_bg_height(n, fsize, avg_chars=avg_ch, area_w_inches=w_in)
 
        bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, BODY_L, BODY_T, left_w, bg_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        _no_shadow(bg)
        _send_to_back(sl, bg)
 
        tb = sl.shapes.add_textbox(BODY_L, BODY_T, left_w, bg_h)
        tf2 = tb.text_frame
        tf2.word_wrap = True
        _set_anchor_middle(tf2)
 
        first = True
        for text in remaining:
            p = tf2.paragraphs[0] if first else tf2.add_paragraph()
            first = False
            _set_bullet_para(p, text, Pt(fsize), DARK_TXT)
 
 
def _add_kpi_dashboard_slide(sl, bullets):
    """
    KPI Dashboard: render bullets as metric tiles rather than bullet points.
    Each bullet is expected to be in the form "Metric: value — context" or similar.
    Falls back to full-width bullets if parsing fails.
    """
    clean = [b for b in (_clean_bullet(x) for x in bullets[:6]) if b]
    if not clean:
        return
 
    # Try to arrange up to 6 tiles in a 3×2 grid
    cols = 3
    rows = (len(clean) + cols - 1) // cols
    tile_w = (BODY_W - Inches(0.3) * (cols - 1)) / cols
    tile_h = (BODY_H - Inches(0.2) * (rows - 1)) / rows
 
    for i, text in enumerate(clean):
        row = i // cols
        col = i % cols
        left = BODY_L + col * (tile_w + Inches(0.3))
        top  = BODY_T + row * (tile_h + Inches(0.2))
 
        tile = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, tile_w, tile_h)
        tile.fill.solid()
        tile.fill.fore_color.rgb = LIGHT_BG
        tile.line.solid()
        tile.line.color.rgb = HP_BLUE
        tile.line.width = Pt(1.5)
        _no_shadow(tile)
 
        tf = tile.text_frame
        tf.word_wrap = True
        tf.margin_left  = Inches(0.12)
        tf.margin_top   = Inches(0.1)
        tf.margin_right = Inches(0.12)
        _set_anchor_middle(tf)
 
        # Try to split "Label: Value — detail" for visual hierarchy
        colon_split = text.split(':', 1)
        if len(colon_split) == 2:
            label_part = colon_split[0].strip()
            value_part = colon_split[1].strip()
 
            p_lbl = tf.paragraphs[0]
            p_lbl.text = label_part.upper()
            p_lbl.font.size = Pt(9)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = HP_BLUE
            p_lbl.alignment = PP_ALIGN.CENTER
 
            p_val = tf.add_paragraph()
            p_val.text = value_part[:60]
            p_val.font.size = Pt(11)
            p_val.font.bold = False
            p_val.font.color.rgb = DARK_TXT
            p_val.alignment = PP_ALIGN.CENTER
        else:
            p = tf.paragraphs[0]
            p.text = text[:80]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TXT
            p.alignment = PP_ALIGN.CENTER
 
 
# ── Chart helpers ─────────────────────────────────────────────────────────────
 
def _chart_pair_title(name1, name2):
    c1, c2 = _clean(name1), _clean(name2)
    words1 = set(re.findall(r'\b[A-Za-z]{3,}\b', c1))
    words2 = set(re.findall(r'\b[A-Za-z]{3,}\b', c2))
    shared = (words1 & words2) - {'Sales','The','And','For','With','Data','Stock','Level'}
    if shared:
        common = sorted(shared, key=len, reverse=True)[0]
        subj1 = re.sub(common, '', c1, flags=re.IGNORECASE).strip(' -–')
        subj2 = re.sub(common, '', c2, flags=re.IGNORECASE).strip(' -–')
        if subj1 and subj2:
            return f"{common}: {subj1} vs {subj2}"
    return f"{c1} & {c2}"[:70]
 
 
def _pair_charts(unplaced):
    remaining = list(unplaced)
    pairs = []
 
    while remaining:
        a = remaining.pop(0)
        if not remaining:
            pairs.append((a,))
            break
 
        akw = _kw(_clean(a[0]))
        best_idx, best_score = 0, -1
        for idx, b in enumerate(remaining):
            score = len(akw & _kw(_clean(b[0])))
            if score > best_score:
                best_score, best_idx = score, idx
 
        b = remaining.pop(best_idx)
        bkw = _kw(_clean(b[0]))
 
        if len(akw) >= 2 and len(bkw) >= 2 and (not (akw - bkw) or not (bkw - akw)):
            pairs.append((a,))
            remaining.insert(0, b)
            continue
 
        pairs.append((a, b, best_score))
 
    return pairs
 
 
# ── Slide builders ────────────────────────────────────────────────────────────
 
def _cover(prs, title):
    sl = _blank(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = HP_BLUE
 
    circ_d = Inches(4.2)
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                SW - circ_d * 0.55, SH - circ_d * 0.55,
                                circ_d, circ_d)
    circ.fill.solid(); circ.fill.fore_color.rgb = HP_RED
    circ.line.fill.background()
    _no_shadow(circ)
    _send_to_back(sl, circ)
 
    tb = sl.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(5), Inches(0.8))
    p  = tb.text_frame.paragraphs[0]
    p.text = "HPGPT"; p.font.size = Pt(28)
    p.font.bold = True; p.font.color.rgb = HP_RED
 
    wrapped = textwrap.fill(_fmt(_flatten(title)), width=34)
    tb2 = sl.shapes.add_textbox(Inches(0.6), Inches(1.8), SW - Inches(1.2), Inches(3.2))
    tf  = tb2.text_frame; tf.word_wrap = True
    p2  = tf.paragraphs[0]
    p2.text = wrapped; p2.font.size = Pt(40)
    p2.font.bold = True; p2.font.color.rgb = WHITE
 
    tb3 = sl.shapes.add_textbox(Inches(0.6), Inches(5.4), SW - Inches(1.2), Inches(0.5))
    p3  = tb3.text_frame.paragraphs[0]
    p3.text = "Generated by HPGPT Agentic System"
    p3.font.size = Pt(13); p3.font.color.rgb = RGBColor(180, 180, 200)
 
    tb4 = sl.shapes.add_textbox(Inches(0.6), Inches(5.95), SW - Inches(1.2), Inches(0.4))
    p4  = tb4.text_frame.paragraphs[0]
    p4.text = datetime.now().strftime("%B %Y")
    p4.font.size = Pt(11); p4.font.color.rgb = RGBColor(160, 160, 180)
 
 
def _thank_you(prs, title):
    sl = _blank(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = HP_BLUE
 
    circ_d = Inches(4.2)
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                -circ_d * 0.55, -circ_d * 0.55,
                                circ_d, circ_d)
    circ.fill.solid(); circ.fill.fore_color.rgb = HP_RED
    circ.line.fill.background()
    _no_shadow(circ)
    _send_to_back(sl, circ)
 
    tb_hp = sl.shapes.add_textbox(Inches(0), SH/2 - Inches(2.0), SW, Inches(1.0))
    p_hp = tb_hp.text_frame.paragraphs[0]
    p_hp.text = "HPCL"; p_hp.font.size = Pt(40)
    p_hp.font.bold = True; p_hp.font.color.rgb = HP_RED
    p_hp.alignment = PP_ALIGN.CENTER
 
    tb2 = sl.shapes.add_textbox(Inches(1), SH/2 - Inches(0.7), SW - Inches(2), Inches(1.0))
    p2  = tb2.text_frame.paragraphs[0]
    p2.text = "Thank You"; p2.font.size = Pt(40)
    p2.font.bold = True; p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
 
    doc_ref  = _flatten(title)[:50]
    date_str = datetime.now().strftime("%d %B %Y")
    tb3 = sl.shapes.add_textbox(Inches(1), SH/2 + Inches(0.4), SW - Inches(2), Inches(0.6))
    p3  = tb3.text_frame.paragraphs[0]
    p3.text = f"HPCL  |  HPGPT  |  {doc_ref}  |  {date_str}"
    p3.font.size = Pt(11); p3.font.color.rgb = RGBColor(180, 180, 200)
    p3.alignment = PP_ALIGN.CENTER
 
 
def _acknowledgement(prs):
    sl = _blank(prs)
 
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, TITLE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = HP_BLUE
    bar.line.fill.background()
    _no_shadow(bar)
 
    tf = bar.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.22)
    p = tf.paragraphs[0]
    p.text = "ACKNOWLEDGEMENT"
    p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT
 
    tb_intro = sl.shapes.add_textbox(BODY_L, BODY_T, BODY_W, Inches(0.45))
    p_intro = tb_intro.text_frame.paragraphs[0]
    p_intro.text = "The author wishes to express sincere gratitude to the following:"
    p_intro.font.size = Pt(13)
    p_intro.font.color.rgb = RGBColor(60, 60, 60)
    p_intro.font.italic = True
 
    placeholders = [
        "[Name], [Designation] — for guidance and mentorship throughout this project.",
        "[Name], [Designation] — for providing access to data and resources.",
        "[Name], [Designation] — for continuous support and valuable feedback.",
        "[Organization / Department] — for the opportunity to undertake this work.",
    ]
    n = len(placeholders)
    bg_h = _estimate_bg_height(n, 15, avg_chars=72)
    tb = sl.shapes.add_textbox(BODY_L, BODY_T + Inches(0.6), BODY_W, bg_h)
    tf2 = tb.text_frame; tf2.word_wrap = True
    _set_anchor_top(tf2)
    first = True
    for ph in placeholders:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        r1 = p.add_run(); r1.text = "• "
        r1.font.size = Pt(15); r1.font.bold = True
        r1.font.color.rgb = HP_RED; r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = ph
        r2.font.size = Pt(15); r2.font.name = "Calibri"
        r2.font.color.rgb = RGBColor(120, 120, 120)
        r2.font.italic = True
        p.space_before = Pt(8); p.space_after = Pt(8)
 
    _footer(sl, "Acknowledgement")
 
 
def _blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl
 
 
def _title_bar(sl, title):
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, TITLE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = HP_BLUE
    bar.line.fill.background()
    _no_shadow(bar)
 
    tf = bar.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.22)
    p  = tf.paragraphs[0]
    p.text = _flatten(title).upper()[:75]
    p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT
 
 
def _footer(sl, doc_title="HPGPT Analytical Report"):
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), FOOTER_T - Inches(0.04), SW, Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = HP_RED
    ln.line.fill.background()
    _no_shadow(ln)
 
    raw = _flatten(doc_title)
    raw = re.sub(r'^\s*HPCL\s*[-:|]?\s*', '', raw, flags=re.IGNORECASE).strip()
    if len(raw) > FOOTER_MAX_CHARS:
        raw = raw[:FOOTER_MAX_CHARS].rsplit(' ', 1)[0] + '…'
    label = f"HPCL  |  {raw}"
    tb = sl.shapes.add_textbox(Inches(0.45), FOOTER_T, Inches(8), Inches(0.38))
    p  = tb.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9); p.font.color.rgb = GREY_TXT
 
    tb2 = sl.shapes.add_textbox(SW - Inches(2.5), FOOTER_T, Inches(2.3), Inches(0.38))
    p2  = tb2.text_frame.paragraphs[0]
    p2.text = datetime.now().strftime("%b %Y")
    p2.font.size = Pt(9); p2.font.color.rgb = GREY_TXT
    p2.alignment = PP_ALIGN.RIGHT
 
 
# ── Bullet area builders ──────────────────────────────────────────────────────
 
def _estimate_bg_height(n_bullets, fsize_pt, avg_chars=60, area_w_inches=12.0):
    chars_per_line = max(1, int(area_w_inches * 9.5))
    wrap_factor    = max(1.0, avg_chars / chars_per_line)
    line_h_inches  = (fsize_pt + _BULLET_PAD) / 72 * wrap_factor
    total          = _TOP_PAD / 72 + n_bullets * line_h_inches + 0.2
    return max(MIN_BG_H, min(Inches(total), BODY_H))
 
 
def _add_bullets_full(sl, bullets):
    clean = [b for b in (_clean_bullet(x) for x in bullets[:MAX_BULLETS]) if b]
    if not clean:
        return
 
    n      = len(clean)
    avg_ch = max((len(b) for b in clean), default=40)
    fsize  = _pick_fsize(n, avg_ch, wide=True)
    w_in   = float(BODY_W / 914400)
    bg_h   = _estimate_bg_height(n, fsize, avg_chars=avg_ch, area_w_inches=w_in)
 
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, BODY_L, BODY_T, BODY_W, bg_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    _no_shadow(bg)
    _send_to_back(sl, bg)
 
    tb = sl.shapes.add_textbox(BODY_L, BODY_T, BODY_W, bg_h)
    tf = tb.text_frame; tf.word_wrap = True
    _set_anchor_middle(tf)
 
    first = True
    for text in clean:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_bullet_para(p, text, Pt(fsize), DARK_TXT)
 
 
def _add_bullets_left(sl, bullets):
    clean = [b for b in (_clean_bullet(x) for x in bullets[:MAX_BULLETS]) if b]
    if not clean:
        return
 
    half_w = SW * 0.52
    n      = len(clean)
    avg_ch = max((len(b) for b in clean), default=40)
    fsize  = _pick_fsize(n, avg_ch, wide=False)
    w_in   = float((half_w - Inches(0.3)) / 914400)
    bg_h   = _estimate_bg_height(n, fsize, avg_chars=avg_ch, area_w_inches=w_in)
 
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, BODY_L, BODY_T, half_w - Inches(0.3), bg_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    _no_shadow(bg)
    _send_to_back(sl, bg)
 
    tb = sl.shapes.add_textbox(BODY_L, BODY_T, half_w - Inches(0.3), bg_h)
    tf = tb.text_frame; tf.word_wrap = True
    _set_anchor_middle(tf)
 
    first = True
    for text in clean:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_bullet_para(p, text, Pt(fsize), DARK_TXT)
 
 
def _pick_fsize(n_bullets, avg_chars, wide=True):
    base = 18 if wide else 15
    if n_bullets <= 2:
        base += 6
    elif n_bullets == 3:
        base += 3
    elif n_bullets >= 7:
        base -= 3
    elif n_bullets >= 5:
        base -= 1
 
    if avg_chars > 120:  base -= 2
    elif avg_chars > 80: base -= 1
    return max(base, 11)
 
 
# ── Chart placement ───────────────────────────────────────────────────────────
 
def _add_chart_right(sl, img_path):
    try:
        from PIL import Image as PI
        img = PI.open(img_path)
        iw, ih = img.size
        area_w = int(SW * 0.44)
        area_h = int(BODY_H)
        scale  = min(area_w/iw, area_h/ih)
        dw, dh = int(iw*scale), int(ih*scale)
        left   = int(SW * 0.54) + (area_w - dw)//2
        top    = int(BODY_T) + (area_h - dh)//2
        sl.shapes.add_picture(img_path, left, top, width=dw, height=dh)
    except Exception as e:
        logger.warning(f"Chart right embed failed: {e}")
 
 
def _add_chart_full(sl, img_path):
    try:
        from PIL import Image as PI
        img = PI.open(img_path)
        iw, ih = img.size
        area_w = int(BODY_W)
        area_h = int(BODY_H)
        scale  = min(area_w/iw, area_h/ih)
        dw, dh = int(iw*scale), int(ih*scale)
        left   = int(BODY_L) + (area_w - dw)//2
        top    = int(BODY_T) + (area_h - dh)//2
        sl.shapes.add_picture(img_path, left, top, width=dw, height=dh)
    except Exception as e:
        logger.warning(f"Chart full embed failed: {e}")
 
 
def _add_two_charts(sl, img1, img2):
    half_w   = int(BODY_W / 2) - int(Inches(0.1))
    area_h   = int(BODY_H)
    area_top = int(BODY_T)
 
    for i, img_path in enumerate([img1, img2]):
        left = int(BODY_L) + i * (half_w + int(Inches(0.2)))
 
        try:
            from PIL import Image as PI
            img = PI.open(img_path)
            iw, ih = img.size
            scale = min(half_w/iw, area_h/ih)
            dw, dh = int(iw*scale), int(ih*scale)
            top  = area_top + (area_h - dh)//2
            img_left = left + (half_w - dw)//2
            sl.shapes.add_picture(img_path, img_left, top, width=dw, height=dh)
        except Exception as e:
            logger.warning(f"Two-chart embed failed (chart {i+1}): {e}")
 
 
# ── Text helpers ──────────────────────────────────────────────────────────────
 
def _set_bullet_para(p, text, fsize, color):
    r1 = p.add_run(); r1.text = "• "
    r1.font.size = fsize; r1.font.bold = True
    r1.font.color.rgb = HP_RED; r1.font.name = "Calibri"
 
    r2 = p.add_run(); r2.text = text
    r2.font.size = fsize; r2.font.name = "Calibri"
    r2.font.color.rgb = color
 
    p.alignment    = PP_ALIGN.LEFT
    p.space_before = Pt(5)
    p.space_after  = Pt(5)
    try:
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(Inches(0.3))))
        pPr.set('indent', str(int(-Inches(0.3))))
    except Exception:
        pass
 
 
def _set_anchor_top(tf):
    try:
        tf._txBody.set(qn('anchor'), 't')
    except Exception:
        pass
 
 
def _set_anchor_middle(tf):
    try:
        tf._txBody.set(qn('anchor'), 'ctr')
    except Exception:
        pass
 
 
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
 
 
def _send_to_back(sl, shape):
    try:
        sp = shape._element
        sl.shapes._spTree.remove(sp)
        sl.shapes._spTree.insert(2, sp)
    except Exception:
        pass
 
 
def _clean_bullet(text):
    text = re.sub(r'^["\u201c\u201d\']+|["\u201c\u201d\']+$', '', _flatten(text)).strip()
    text = re.sub(r'["\u201c\u201d]([^"]+)["\u201c\u201d]', r'\1', text)
    text = re.sub(r'^[-•*+\d.)]+\s*', '', text).strip()
    text = re.sub(r'[*_]+', '', text)
 
    orphan_patterns = [
        r'can be observed from the following chart',
        r'is presented in the following chart',
        r'can be visualized as follows',
        r'is shown in the (following |)chart',
        r'as shown in the (following |)chart',
        r'the (sales |revenue |data |trend )?(trend|chart|graph|figure) (below|above|following)',
        r'the following (charts?|graphs?|figures?) illustrate',
        r'the following (charts?|graphs?|figures?) (shows?|show|presents?|present)',
        r'chart candidate',
        r'numeric column',
        r'text column',
        r'the data shows that the data',
        r'refer to the (chart|graph|figure)',
        r'see (the )?(chart|graph|figure) (below|above)',
    ]
    for pat in orphan_patterns:
        if re.search(pat, text, flags=re.IGNORECASE):
            return ''
 
    return text.strip()
 
 
# ── Content parsing ───────────────────────────────────────────────────────────
 
_SENTENCE_END = re.compile(r'(?<=[.!?])\s+')
 
# Slide type detection from heading text
_SLIDE_TYPE_MAP = [
    (re.compile(r'\bexecutive\s+summary\b', re.IGNORECASE), "executive_summary"),
    (re.compile(r'\bkpi\s+(dashboard|overview|summary)\b', re.IGNORECASE), "kpi_dashboard"),
    (re.compile(r'\bperformance\s+(dashboard|scorecard|overview)\b', re.IGNORECASE), "kpi_dashboard"),
]
 
 
def _detect_slide_type(heading: str) -> str:
    for pattern, stype in _SLIDE_TYPE_MAP:
        if pattern.search(heading):
            return stype
    return "content"
 
 
def _split_sentences(text):
    parts = _SENTENCE_END.split(text.strip())
    result, current = [], ""
    for part in parts:
        candidate = f"{current} {part}".strip() if current else part
        if len(candidate) > 160 and current:
            result.append(current)
            current = part
        else:
            current = candidate
    if current:
        result.append(current)
    return result or [text]
 
 
def _parse(content):
    """
    Parse LLM markdown into (title, bullets, slide_type) triples.
    slide_type is detected from the heading for layout-aware rendering.
    """
    global _seen_headings
    lines = content.splitlines()
    raw_sections = []
    cur_title, cur_bullets, cur_type = "Overview", [], "content"
 
    for line in lines:
        s = line.strip()
 
        if re.match(r'^\[CHART:', s, re.IGNORECASE):
            continue
 
        m = re.match(r'^(#{1,3})\s+(.+)$', s)
        if m:
            if cur_bullets and cur_title is not None:
                raw_sections.append((cur_title, cur_bullets[:], cur_type))
            heading = re.sub(r'[*_]+', '', m.group(2)).strip()
 
            if _FILTERED_HEADINGS.match(heading):
                cur_title = None
                cur_bullets = []
                cur_type = "content"
                continue
 
            heading_key = re.sub(r'\s+', ' ', heading.lower().strip())
            if heading_key in _seen_headings:
                cur_title = None
                cur_bullets = []
                cur_type = "content"
                continue
            _seen_headings.add(heading_key)
 
            cur_title   = heading
            cur_bullets = []
            cur_type    = _detect_slide_type(heading)
            continue
 
        if cur_title is None: continue
        if '|' in s:          continue
        if not s:              continue
 
        text = re.sub(r'^[-•*+\d.)]+\s*', '', s)
        text = re.sub(r'[*_]+', '', text).strip()
        if len(text) < 8:     continue
 
        for sent in _split_sentences(text):
            sent = sent.strip()
            if sent:
                cur_bullets.append(sent)
 
    if cur_bullets and cur_title is not None:
        raw_sections.append((cur_title, cur_bullets, cur_type))
 
    # ── Merge small sections ──────────────────────────────────────────────────
    merged, i = [], 0
    while i < len(raw_sections):
        title, bullets, stype = raw_sections[i]
        while (len(bullets) < 4 and i + 1 < len(raw_sections)
               and len(merged) < MAX_SLIDES):
            nt, nb, ntype = raw_sections[i+1]
            if stype == "content" and ntype == "content" and len(bullets) + len(nb) <= MAX_BULLETS:
                bullets = bullets + nb
                i += 1
            else:
                break
        if len(bullets) > MAX_BULLETS:
            mid = min((len(bullets) + 1) // 2, MAX_BULLETS)
            if len(bullets) - mid < 2:
                mid = max(2, len(bullets) - 2)
            merged.append((title, bullets[:mid], stype))
            merged.append((title + " (cont.)", bullets[mid:], "content"))
        else:
            merged.append((title, bullets, stype))
        i += 1
 
    return merged
 
 
# ── Utilities ─────────────────────────────────────────────────────────────────
 
STOP = {'by','and','the','of','in','a','an','to','for','from','with',
        'data','analysis','report','file','product','wise','per','vs',
        'are','has','its','this','that','these','all','each'}
 
def _kw(text):
    return set(re.findall(r'[a-z]{3,}', text.lower())) - STOP
 
def _clean(name):
    name = re.sub(r'—.*$', '', name).strip()
    name = name.replace('_KL','').replace('_MT','').replace('_INR','')
    parts = name.replace('_',' ').split()
    return ' '.join(p.upper() if p.upper() in ABBREVS else p.title() for p in parts)
 
def _flatten(item):
    if isinstance(item, (list, tuple)):
        return ' '.join(_flatten(i) for i in item)
    return str(item)
 
def _fmt(t):
    skip = {'and','or','the','of','in','on','with','a','an','to','for'}
    words = t.strip().split()
    result = []
    for i, w in enumerate(words):
        if w.upper() in ABBREVS:
            result.append(w.upper())
        elif i == 0 or i == len(words)-1 or w.lower() not in skip:
            result.append(w.capitalize())
        else:
            result.append(w.lower())
    return ' '.join(result)
 
def wants_acknowledgement(prompt: str) -> bool:
    return bool(re.search(
        r'\b(acknowledgement|acknowledgment|acknowledge)\b',
        prompt, re.IGNORECASE
    ))